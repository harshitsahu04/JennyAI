import streamlit as st
import os, io, re, zipfile, glob, requests, json
from dataclasses import dataclass
from typing import List, Tuple, Optional
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font
from docx import Document
from docx.shared import Pt as DOCXPt
from docx.enum.text import WD_ALIGN_PARAGRAPH

# ==============================
# --- PREMIUM CSS (BCG-style) ---
# ==============================
st.markdown("""
<style>
.stApp {background-color: #F5F5F2; color: #1A1A1A; font-family: 'Helvetica Neue', sans-serif;}
.hero {background: linear-gradient(120deg, #0A2342, #1C335D); border-radius:12px; padding:30px; margin-bottom:30px; text-align:center; color:#FFD700; box-shadow:0 6px 20px rgba(0,0,0,0.12);}
.hero h1 {font-size:2.6em; margin-bottom:8px; color:#CBCBD4; font-weight:800; letter-spacing:1px;}
.hero p {font-size:1.1em; color:#FFD700; font-weight:500; margin-top:0;}
h1 {color:#0A2342; font-weight:700; text-align:center; font-size:2.4em; border-bottom:3px solid #D4AF37; padding-bottom:10px; margin-bottom:25px;}
h3 {color:#5A5F6C; text-align:center; margin-top:-10px; margin-bottom:40px; font-weight:400;}
label {color:#0A2342 !important; opacity:1 !important; font-weight:700; font-size:16px; margin-bottom:6px; display:block; position:relative; padding-left:12px;}
label:before {content:""; position:absolute; left:0; top:50%; transform:translateY(-50%); width:4px; height:16px; background-color:#D4AF37; border-radius:2px;}
.stTextArea textarea {font-family:"Georgia", serif; background-color:#FFF; border:1px solid #D4AF37; border-radius:12px; font-size:15px; padding:14px; color:#0A2342; caret-color:#D4AF37; box-shadow:0 2px 6px rgba(0,0,0,0.08);}
.stTextArea textarea:focus {border-color:#B8860B; box-shadow:0 0 10px rgba(212,175,55,0.25);}
div.stButton > button:first-child {background:linear-gradient(135deg,#1C335D,#0A2342); color:#FFD700; border-radius:12px; font-weight:700; font-size:16px; padding:14px 32px; box-shadow:0 4px 14px rgba(0,0,0,0.1); display:block; margin:0 auto; margin-top:20px; transition: all 0.3s ease;}
div.stButton > button:first-child:hover {background:linear-gradient(135deg,#0A2342,#1C335D); transform:translateY(-2px); box-shadow:0 6px 18px rgba(0,0,0,0.15);}
.response-box {background:linear-gradient(90deg,#FFFDF7,#F7F5EE); border-left:6px solid #D4AF37; padding:28px; border-radius:12px; margin-top:30px; font-size:15px; color:#0A2342; box-shadow:0 6px 20px rgba(0,0,0,0.1);}
.export-title {font-size:18px; font-weight:600; color:#0A2342; margin-top:35px; margin-bottom:10px; text-align:center;}
.stAlert {background-color:#FFEAEA !important; border:1px solid #E0B4B4 !important; border-radius:8px;}
.stAlert p, .stAlert span {color:#5A2E2E !important; font-weight:500;}
</style>
""", unsafe_allow_html=True)

# ==============================
# --- HEADER ---
# ==============================
st.markdown('<div class="hero"><h1>JENNY</h1><p>Your Elite Strategy Partner</p></div>', unsafe_allow_html=True)

# ==============================
# --- INPUTS ---
# ==============================
problem = st.text_area("Problem Statement", height=120)
context = st.text_area("Client Context", height=120)

# ==============================
# --- CORPUS & TF-IDF ---
# ==============================
PDF_DIR = "jenny pdfs"  # change to your PDF folder
os.makedirs(PDF_DIR, exist_ok=True)

@dataclass
class CorpusIndex:
    docs: List[str]
    vectorizer: Optional[TfidfVectorizer]
    matrix: Optional[np.ndarray]
    sources: List[Tuple[str, int]]

def extract_text_from_pdf("jenny pdfs/": str) -> str:
    import fitz
    parts = []
    try:
        with fitz.open("jenny pdfs/") as doc:
            for page in doc:
                parts.append(page.get_text("text") or "")
    except:
        return ""
    return "\n".join(parts)

def chunk_text(text: str, max_chars: int = 1800, overlap: int = 220) -> List[str]:
    text = re.sub(r"\s+", " ", text or "").strip()
    chunks, i = [], 0
    while i < len(text):
        end = min(i + max_chars, len(text))
        chunks.append(text[i:end])
        i = end - overlap if end - overlap > i else end
    return chunks

def build_index(PDF_DIR: str) -> CorpusIndex:
    pdf_path = sorted(glob.glob(os.path.join(PDF_DIR, "**/*.pdf"), recursive=True))
    all_chunks, sources = [], []
    for p in pdf_path:
        raw = extract_text_from_pdf(p)
        for idx, ch in enumerate(chunk_text(raw)):
            all_chunks.append(ch)
            sources.append((os.path.basename(p), idx))
    if not all_chunks:
        all_chunks, sources = ["(No docs indexed)"], [("—", 0)]
    vec = TfidfVectorizer(min_df=1, max_df=0.95)
    mat = vec.fit_transform(all_chunks)
    return CorpusIndex(all_chunks, vec, mat, sources)

# ✅ FIXED RETRIEVE FUNCTION
def retrieve(corpus: CorpusIndex, query: str, k: int = 5):
    if corpus is None or corpus.vectorizer is None or corpus.matrix is None:
        return []
    try:
        sims = cosine_similarity(corpus.vectorizer.transform([query]), corpus.matrix).ravel()
        topk = sims.argsort()[::-1][:k]
        return [(corpus.docs[i], float(sims[i])) for i in topk]
    except Exception as e:
        print("⚠️ Error in retrieve:", e)
        return []

corpus = build_index(PDF_DIR)

# ==============================
# --- GROQ REST HELPER ---
# ==============================
GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "").strip()
GROQ_API_URL = "https://api.groq.com/openai/v1/chat/completions"

SYSTEM_PROMPT = "You are Jenny, a senior strategy consultant. STRICTLY business/strategy only."

def jenny_answer(question: str, corpus):
    hits = retrieve(corpus, question, k=5)
    ctx = "\n".join([h[0] for h in hits])[:2400] if hits else "(no context)"
    payload = {
        "model": "llama-3.1-8b-instant",
        "messages": [
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": f"{question}\nContext:\n{ctx}"}
        ],
        "temperature": 0.2,
        "max_tokens": 2000
    }
    try:
        resp = requests.post(
            GROQ_API_URL,
            headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
            data=json.dumps(payload),
            timeout=60
        )
        if resp.status_code == 200:
            return resp.json()["choices"][0]["message"]["content"]
    except Exception as e:
        print("⚠️ Groq REST call failed:", e)
    return "(Jenny could not fetch a live response; fallback used.)"

# ==============================
# --- EXPORT FUNCTIONS ---
# ==============================
def make_pptx(ans, q):  # simple 2-slide for demo
    prs = Presentation()
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    s.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8), Inches(1)).text_frame.add_paragraph().text = q
    s = prs.slides.add_slide(blank)
    s.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8), Inches(5)).text_frame.add_paragraph().text = ans
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()

def make_docx(ans, q):
    doc = Document()
    doc.add_heading(q, 0)
    doc.add_paragraph(ans)
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.read()

def make_excel():
    wb = Workbook()
    ws = wb.active
    ws.title = "Demo"
    ws.append(["Metric", "Value"])
    ws.append(["ROI", 0.25])
    bio = io.BytesIO()
    wb.save(bio)
    bio.seek(0)
    return bio.read()

# ==============================
# --- MAIN BUTTON LOGIC ---
# ==============================
if st.button("Ask Jenny"):
    if not problem.strip():
        st.error("Enter Problem Statement")
        st.stop()
    with st.spinner("Jenny is thinking..."):
        answer = jenny_answer(problem, corpus)
    st.markdown(f'<div class="response-box">{answer}</div>', unsafe_allow_html=True)

    st.markdown('<div class="export-title">📤 Export Results</div>', unsafe_allow_html=True)
    ppt_bytes, doc_bytes, xls_bytes = make_pptx(answer, problem), make_docx(answer, problem), make_excel()
    col1, col2, col3 = st.columns(3)
    with col1:
        st.download_button("📊 PPTX", ppt_bytes, file_name="jenny.pptx")
    with col2:
        st.download_button("📄 DOCX", doc_bytes, file_name="jenny.docx")
    with col3:
        st.download_button("📈 XLSX", xls_bytes, file_name="jenny.xlsx")
